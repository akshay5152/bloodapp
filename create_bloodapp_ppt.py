#!/usr/bin/env python3
"""
BloodApp PowerPoint Presentation Generator
Creates a professional PPT file for BloodApp architecture presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_bloodapp_presentation():
    """Create BloodApp PowerPoint presentation"""
    
    # Create presentation
    prs = Presentation()
    
    # Define colors
    primary_color = RGBColor(102, 126, 234)  # #667eea
    secondary_color = RGBColor(118, 75, 162)  # #764ba2
    dark_text = RGBColor(44, 62, 80)  # #2c3e50
    light_text = RGBColor(127, 140, 141)  # #7f8c8d
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.title
    title.text = "BloodApp"
    title.text_frame.paragraphs[0].font.size = Pt(48)
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    title.text_frame.paragraphs[0].font.bold = True
    
    # Subtitle
    subtitle = slide.placeholders[1]
    subtitle.text = "Blood Donation Management System"
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.color.rgb = light_text
    
    # Add content box
    content_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(3))
    content_frame = content_box.text_frame
    content_frame.text = "🏥 Modern Blood Donation Platform\n\nA comprehensive Progressive Web App (PWA) designed to connect donors, recipients, and hospitals in real-time blood donation management."
    content_frame.paragraphs[0].font.size = Pt(20)
    content_frame.paragraphs[0].font.color.rgb = primary_color
    content_frame.paragraphs[0].font.bold = True
    content_frame.paragraphs[2].font.size = Pt(14)
    content_frame.paragraphs[2].font.color.rgb = dark_text
    
    # Slide 2: System Architecture
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "System Architecture"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    # Add architecture diagram
    content = slide.placeholders[1]
    content_frame = content.text_frame
    content_frame.text = "Frontend Layer\n• React 19 + TypeScript + React Router\n\nState Management\n• React Context API + Local Storage\n\nUI Framework\n• Custom CSS + Responsive Design\n\nPWA Features\n• Service Worker + Offline Support"
    
    # Format the content
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = dark_text
        if paragraph.text.startswith("Frontend Layer") or paragraph.text.startswith("State Management") or paragraph.text.startswith("UI Framework") or paragraph.text.startswith("PWA Features"):
            paragraph.font.bold = True
            paragraph.font.color.rgb = primary_color
    
    # Slide 3: Technology Stack
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Technology Stack"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide.placeholders[1]
    content_frame = content.text_frame
    content_frame.text = "React v19.1.0 - Frontend Framework\nTypeScript v4.9.5 - Type Safety\nReact Router v7.7.1 - Client-side Routing\nChart.js v4.5.0 - Data Visualization\nCapacitor v7.4.2 - Mobile App Conversion\nPWA Native - Progressive Web App"
    
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = dark_text
    
    # Slide 4: User Types & Features
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "User Types & Features"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide.placeholders[1]
    content_frame = content.text_frame
    content_frame.text = "👤 Regular Users\n• OTP-based authentication\n• Blood donation requests\n• Profile management\n• Request history\n• Dashboard analytics\n\n🏥 Hospitals\n• Hospital registration\n• Blood inventory management\n• Donor coordination\n• Emergency requests\n• Event management\n\n👨‍💼 Admin Panel\n• User management\n• System analytics\n• Emergency broadcasts\n• Event invitations\n• Blood group requests"
    
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = dark_text
        if paragraph.text.startswith("👤") or paragraph.text.startswith("🏥") or paragraph.text.startswith("👨‍💼"):
            paragraph.font.bold = True
            paragraph.font.color.rgb = primary_color
    
    # Slide 5: Application Flow
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Application Flow"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide.placeholders[1]
    content_frame = content.text_frame
    content_frame.text = "1. Authentication\n   User login with mobile OTP or admin credentials\n\n2. Dashboard Access\n   Role-based dashboard with analytics and quick actions\n\n3. Blood Requests\n   Create and manage blood donation requests\n\n4. Matching & Coordination\n   System matches donors with requests\n\n5. Emergency Management\n   Admin can send emergency broadcasts to users/hospitals"
    
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = dark_text
        if paragraph.text.startswith("1.") or paragraph.text.startswith("2.") or paragraph.text.startswith("3.") or paragraph.text.startswith("4.") or paragraph.text.startswith("5."):
            paragraph.font.bold = True
            paragraph.font.color.rgb = primary_color
    
    # Slide 6: Key Features
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Key Features"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide.placeholders[1]
    content_frame = content.text_frame
    content_frame.text = "🚨 Emergency Requests\nAdmin can send emergency blood group requests or event invitations to users/hospitals via email or SMS\n\n📊 Analytics Dashboard\nReal-time charts and statistics for blood donation trends and user activity\n\n📱 Mobile-First Design\nResponsive design optimized for mobile devices with PWA capabilities\n\n🔒 Secure Storage\nLocal storage for session management and offline data persistence\n\n🎨 Theme Management\nLight theme configuration with modern UI/UX design patterns\n\n⚡ Performance Optimized\nFast loading with optimized bundles and lazy loading components"
    
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = dark_text
        if paragraph.text.startswith("🚨") or paragraph.text.startswith("📊") or paragraph.text.startswith("📱") or paragraph.text.startswith("🔒") or paragraph.text.startswith("🎨") or paragraph.text.startswith("⚡"):
            paragraph.font.bold = True
            paragraph.font.color.rgb = primary_color
    
    # Slide 7: Deployment & Scalability
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Deployment & Scalability"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide.placeholders[1]
    content_frame = content.text_frame
    content_frame.text = "☁️ Cloud Deployment\nReady for deployment on Vercel, Netlify, Firebase, or GitHub Pages\n\n📱 Mobile App Conversion\nCapacitor integration for native mobile app conversion\n\n🔧 Easy Maintenance\nComponent-based architecture for easy updates and feature additions\n\n🚀 Ready for Production\n• Optimized build process with React Scripts\n• PWA configuration for mobile deployment\n• TypeScript for better code quality\n• Responsive design for all devices"
    
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = dark_text
        if paragraph.text.startswith("☁️") or paragraph.text.startswith("📱") or paragraph.text.startswith("🔧") or paragraph.text.startswith("🚀"):
            paragraph.font.bold = True
            paragraph.font.color.rgb = primary_color
    
    # Slide 8: Future Roadmap
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Future Roadmap"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide.placeholders[1]
    content_frame = content.text_frame
    content_frame.text = "🔗 Backend Integration\nConnect to Node.js/Express or Firebase backend for real-time data\n\n🔔 Push Notifications\nReal-time notifications for blood requests and updates\n\n🗺️ Location Services\nGPS integration for nearby donor/hospital matching\n\n📈 Advanced Analytics\nMachine learning for predictive blood demand"
    
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = dark_text
        if paragraph.text.startswith("🔗") or paragraph.text.startswith("🔔") or paragraph.text.startswith("🗺️") or paragraph.text.startswith("📈"):
            paragraph.font.bold = True
            paragraph.font.color.rgb = primary_color
    
    # Save the presentation
    output_file = "BloodApp_Architecture_Presentation.pptx"
    prs.save(output_file)
    
    print(f"✅ PowerPoint presentation created successfully!")
    print(f"📁 File saved as: {output_file}")
    print(f"📍 Location: {os.path.abspath(output_file)}")
    
    return output_file

if __name__ == "__main__":
    try:
        # Check if python-pptx is installed
        import pptx
        create_bloodapp_presentation()
    except ImportError:
        print("❌ Error: python-pptx library not found!")
        print("📦 Please install it using: pip install python-pptx")
        print("\n🔧 Installation steps:")
        print("1. Open terminal/command prompt")
        print("2. Run: pip install python-pptx")
        print("3. Run this script again")
    except Exception as e:
        print(f"❌ Error creating presentation: {e}") 